"""
ConsultMatch FastAPI backend — v13
Role-level postings. Consultants apply to roles, not projects.
"""
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
import os, httpx

from data.synthetic import CONSULTANTS, PROJECTS, ROLES
from scoring import compute_score, rank_roles_for_consultant, rank_consultants_for_role

app = FastAPI(title="ConsultMatch API", version="3.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "https://consultmatch.vercel.app",
        "https://consultmatch-xi.vercel.app",
    ],
    allow_origin_regex=r"https://.*\.vercel\.app",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── In-memory state ───────────────────────────────────────────
_custom_consultants: dict[str, dict] = {}
_applications: dict[str, list] = {}   # consultant_id -> [{role_id, status, cv_text, ...}]
_likes: dict[str, list] = {}           # consultant_id -> [role_id, ...]

_p_lookup = {p["id"]: p for p in PROJECTS}
_r_lookup = {r["role_id"]: r for r in ROLES}

def get_all_consultants():
    return CONSULTANTS + list(_custom_consultants.values())

def get_c_lookup():
    return {c["id"]: c for c in get_all_consultants()}


# ── Models ────────────────────────────────────────────────────
class CustomProfile(BaseModel):
    id: str; name: str; cl: int; cl_title: str
    skills: list[str]; preferred_industries: list[str]
    wfh_preference: str; work_style: str
    preferred_duration: str; career_goal: str
    available_from: str
    cvText: Optional[str] = ""
    cvFileName: Optional[str] = ""
    cvFileBase64: Optional[str] = ""

class ApplicationRequest(BaseModel):
    consultant_id: str; role_id: str
    cv_text: Optional[str] = ""
    # Optional metadata for custom project slots
    role_title: Optional[str] = None
    project_name: Optional[str] = None
    client: Optional[str] = None
    industry: Optional[str] = None
    district: Optional[str] = None
    cl_label: Optional[str] = None

class LikeRequest(BaseModel):
    consultant_id: str; role_id: str

class CustomScoreRequest(BaseModel):
    consultant: dict; role_id: str

class TailorRequest(BaseModel):
    cv_text: Optional[str] = ""
    role_id: str
    file_base64: Optional[str] = None
    file_name: Optional[str] = None

class StatusUpdate(BaseModel):
    consultant_id: str; role_id: str; status: str


# ── Custom profile ────────────────────────────────────────────
@app.post("/consultants/register")
def register_custom_consultant(profile: CustomProfile):
    _custom_consultants[profile.id] = profile.model_dump()
    return {"status": "registered", "id": profile.id}


# ── Consultants ───────────────────────────────────────────────
@app.get("/consultants")
def list_consultants():
    return get_all_consultants()

@app.get("/consultants/{consultant_id}")
def get_consultant(consultant_id: str):
    c = get_c_lookup().get(consultant_id)
    if not c: raise HTTPException(404, "Not found")
    return c

@app.get("/consultants/{consultant_id}/recommendations")
def get_recommendations(consultant_id: str):
    c = get_c_lookup().get(consultant_id)
    if not c: raise HTTPException(404, "Not found")
    c_cl = c.get("cl", 9)
    eligible_roles = [r for r in ROLES if r["cl_range"][0] <= c_cl <= r["cl_range"][1]]
    return rank_roles_for_consultant(c, eligible_roles)


# ── Roles ─────────────────────────────────────────────────────
@app.get("/roles")
def list_roles():
    return ROLES

@app.get("/roles/{role_id}")
def get_role(role_id: str):
    r = _r_lookup.get(role_id)
    if not r: raise HTTPException(404, "Not found")
    return r

@app.get("/roles/{role_id}/candidates")
def get_role_candidates(role_id: str):
    r = _r_lookup.get(role_id)
    if not r: raise HTTPException(404, "Not found")
    applicant_by_cid = {
        cid: a for cid, apps in _applications.items()
        for a in apps if a["role_id"] == role_id
    }
    liked_by = {cid for cid, rids in _likes.items() if role_id in rids}
    candidates = []
    for c in rank_consultants_for_role(r, get_all_consultants()):
        cid = c["id"]
        candidates.append({
            **c,
            "has_applied": cid in applicant_by_cid,
            "has_liked": cid in liked_by,
            "cv_text": applicant_by_cid.get(cid, {}).get("cv_text",""),
            "application_status": applicant_by_cid.get(cid, {}).get("status",""),
        })
    return candidates


# ── Projects ──────────────────────────────────────────────────
@app.get("/projects")
def list_projects():
    return PROJECTS

@app.get("/projects/{project_id}")
def get_project(project_id: str):
    p = _p_lookup.get(project_id)
    if not p: raise HTTPException(404, "Not found")
    return p

@app.get("/projects/{project_id}/roles")
def get_project_roles(project_id: str):
    """All role postings for a specific project."""
    return [r for r in ROLES if r["project_id"] == project_id]

@app.post("/projects/{project_id}/propose-team")
def propose_team(project_id: str):
    """Propose best team for a project across all its role slots."""
    project_roles = [r for r in ROLES if r["project_id"] == project_id]
    if not project_roles: raise HTTPException(404, "No roles found")

    all_consultants = get_all_consultants()
    applicant_by_cid = {
        cid: a for cid, apps in _applications.items()
        for a in apps if a.get("role_id","").startswith(project_id)
    }
    liked_by = {cid for cid, rids in _likes.items()
                if any(rid.startswith(project_id) for rid in rids)}

    proposed = []
    used = set()
    for role in project_roles:
        candidates = rank_consultants_for_role(role, all_consultants)
        scored = []
        for c in candidates:
            if c["id"] in used: continue
            cid = c["id"]
            boost = 10 if cid in applicant_by_cid else (5 if cid in liked_by else 0)
            scored.append({**c, "adjusted": c["score"]["total"] + boost,
                          "has_applied": cid in applicant_by_cid,
                          "has_liked": cid in liked_by})
        scored.sort(key=lambda x: x["adjusted"], reverse=True)
        best = scored[0] if scored else None
        if best: used.add(best["id"])
        proposed.append({"role": role, "proposed": [best] if best else []})

    return {"project_id": project_id, "proposed_team": proposed}


# ── Likes ─────────────────────────────────────────────────────
@app.post("/likes")
def toggle_like(req: LikeRequest):
    liked = _likes.setdefault(req.consultant_id, [])
    if req.role_id in liked:
        liked.remove(req.role_id)
        return {"status": "unliked"}
    liked.append(req.role_id)
    return {"status": "liked"}

@app.get("/likes/{consultant_id}")
def get_likes(consultant_id: str):
    return {"liked": _likes.get(consultant_id, [])}


# ── Applications ──────────────────────────────────────────────
@app.post("/applications")
def apply_to_role(req: ApplicationRequest):
    apps = _applications.setdefault(req.consultant_id, [])
    if any(a["role_id"] == req.role_id for a in apps):
        raise HTTPException(400, "Already applied to this role")
    from datetime import datetime
    apps.append({
        "role_id": req.role_id,
        "status": "Applied",
        "cv_text": req.cv_text,
        "applied_at": datetime.utcnow().isoformat(),
        # Store metadata for custom slots
        "role_title": req.role_title,
        "project_name": req.project_name,
        "client": req.client,
        "industry": req.industry,
        "district": req.district,
        "cl_label": req.cl_label,
    })
    return {"status": "ok"}

@app.get("/applications/{consultant_id}")
def get_applications(consultant_id: str):
    apps = _applications.get(consultant_id, [])
    enriched = []
    for a in apps:
        r = _r_lookup.get(a["role_id"])
        if r:
            enriched.append({**a,
                "role_title": r["role_title"],
                "project_name": r["project_name"],
                "client": r["client"],
                "industry": r["industry"],
                "district": r.get("district","Singapore"),
                "cl_label": r["cl_label"],
            })
        else:
            # Custom project slot — show with stored info
            enriched.append({**a,
                "role_title": a.get("role_title", "Role"),
                "project_name": a.get("project_name", "Custom Project"),
                "client": a.get("client", "—"),
                "industry": a.get("industry", "—"),
                "district": a.get("district", "Singapore"),
                "cl_label": a.get("cl_label", "—"),
            })
    return enriched

@app.patch("/applications/status")
def update_application_status(req: StatusUpdate):
    apps = _applications.get(req.consultant_id, [])
    for a in apps:
        if a["role_id"] == req.role_id:
            a["status"] = req.status
            return {"status": "updated"}
    raise HTTPException(404, "Application not found")


# ── Scoring ───────────────────────────────────────────────────
@app.post("/score/custom")
def get_custom_score(request: CustomScoreRequest):
    r = _r_lookup.get(request.role_id)
    if not r: raise HTTPException(404, "Role not found")
    return compute_score(request.consultant, r)


# ── CV Tailoring ──────────────────────────────────────────────
def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    import io
    name = filename.lower()
    if name.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    elif name.endswith(".docx") or name.endswith(".doc"):
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ""

@app.post("/tailor")
async def tailor_cv(request: TailorRequest):
    r = _r_lookup.get(request.role_id)
    if not r: raise HTTPException(404, "Role not found")
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key: raise HTTPException(500, "ANTHROPIC_API_KEY not configured")

    role_context = (
        f"Role: {r['role_title']} — {r['project_name']} ({r['client']}, {r['industry']})"
        f"\nRequired skills: {', '.join(r['required_skills'])}"
        f"\nJob description: {r['description']}"
    )

    tailor_instructions = (
        "You are a senior CV writer at a top management consulting firm, expert at tailoring CVs for specific roles.\n\n"
        f"The consultant is applying to:\n{role_context}\n\n"
        "Your task is to aggressively tailor the experience section for this specific role. Follow these rules:\n\n"
        "1. REORDER: Within each job, move bullets most relevant to the target role to the top. Bury or remove irrelevant bullets.\n"
        "2. REFRAME: Rephrase bullets to use language from the job description and highlight transferable skills. Name exact tools and methods that match required skills.\n"
        "3. QUANTIFY: Where numbers exist in the original, keep them. Add specificity where the original is vague.\n"
        "4. DOWNPLAY: Shrink or remove bullets that add no value for this role. Dropping bullets entirely is acceptable if they are irrelevant.\n"
        "5. KEEP FACTS: Do not invent experience, numbers, or technologies not in the original CV.\n"
        "6. KEEP STRUCTURE: Preserve all job titles, company names, and date ranges exactly as-is.\n"
        "7. OUTPUT FORMAT: Same format as input (bullet points under each job). Output ONLY the rewritten experience section — no preamble, no explanation.\n\n"
        "Be bold. A generic CV that slightly reorders bullets is not acceptable. The output should read as if this consultant was born for this specific role."
    )

    if request.file_base64 and request.file_name:
        import base64
        file_bytes = base64.b64decode(request.file_base64)
        name = request.file_name.lower()
        if name.endswith(".pdf"):
            messages = [{"role": "user", "content": [
                {"type": "document", "source": {
                    "type": "base64", "media_type": "application/pdf",
                    "data": request.file_base64
                }},
                {"type": "text", "text": tailor_instructions + "\n\nThe CV is attached as a PDF. Extract the experience section and tailor it as instructed."}
            ]}]
        else:
            extracted = extract_text_from_file(file_bytes, request.file_name)
            if not extracted:
                raise HTTPException(400, "Could not extract text from file")
            messages = [{"role": "user", "content":
                tailor_instructions + f"\n\nHere is the CV content extracted from {request.file_name}:\n\n{extracted}"}]
    else:
        messages = [{"role": "user", "content":
            tailor_instructions + f"\n\nHere is the consultant's experience section:\n\n{request.cv_text}"}]

    async with httpx.AsyncClient(timeout=60.0) as client:
        response = await client.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json"
            },
            json={"model": "claude-sonnet-4-6", "max_tokens": 2000, "messages": messages},
        )
    if response.status_code != 200:
        raise HTTPException(502, f"Anthropic API error: {response.text}")
    return {"tailored_text": response.json()["content"][0]["text"], "mode": "ai"}


# ── Reset ─────────────────────────────────────────────────────
@app.post("/reset")
def reset_state():
    _custom_consultants.clear()
    _applications.clear()
    _likes.clear()
    return {"status": "reset"}


# ── Batch scoring ─────────────────────────────────────────────
class BatchScoreRequest(BaseModel):
    consultant: dict

@app.post("/score/custom/batch")
def batch_score_custom(request: BatchScoreRequest):
    """Score a custom consultant against eligible roles (CL-filtered) in one call."""
    c_cl = request.consultant.get("cl", 9)
    eligible_roles = [r for r in ROLES if r["cl_range"][0] <= c_cl <= r["cl_range"][1]]
    results = []
    for role in eligible_roles:
        score = compute_score(request.consultant, role)
        results.append({**role, "score": score})
    results.sort(key=lambda x: x["score"]["total"], reverse=True)
    return results


class InlineScoreRequest(BaseModel):
    consultant: dict  # full consultant object
    role: dict        # full role definition

@app.post("/score/inline")
def score_inline(request: InlineScoreRequest):
    """Score a consultant against an inline role definition (for custom projects)."""
    return compute_score(request.consultant, request.role)
